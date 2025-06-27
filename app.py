import os
import re
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from PIL import Image
from google import genai
import pytesseract
import random
import io
import platform
import shutil
from flask import Flask
from flask import Flask,jsonify
from queue import Queue
from threading import Thread, Semaphore,Lock
import database as db
import requests
import psutil
import gc
# import tempfile
from contextlib import contextmanager

# Set temp directory path early so it's available everywhere
TEMP_DIR = os.path.join(os.path.dirname(__file__), 'temp') if platform.system() == 'Windows' else '/tmp/notesup_temp'

queue = Queue()
semaphore = Semaphore(1)

worker_active = False
worker_lock = Lock()

app = Flask(__name__)

# Configure Tesseract path based on environment
def configure_tesseract():
    """Configure Tesseract OCR path based on the operating system"""
    if platform.system() == "Windows":
        # Windows path (for local development)
        tesseract_path = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
        if os.path.exists(tesseract_path):
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
        else:
            print("Warning: Tesseract not found at expected Windows path")
    elif platform.system() == "Linux":
        # Linux path (for Docker container)
        tesseract_path = shutil.which('tesseract')
        if tesseract_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
        else:
            print("Warning: Tesseract not found in PATH")
    else:
        # macOS or other systems
        tesseract_path = shutil.which('tesseract')
        if tesseract_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path

# Initialize Tesseract configuration
configure_tesseract()

def clean_text(text):
    return re.sub(r'\s+', ' ', text).strip()

def log_memory_usage(context=""):
    process = psutil.Process(os.getpid())
    mem_mb = process.memory_info().rss / (1024 * 1024)
    print(f"[MEMORY] {context} - RSS: {mem_mb:.2f} MB")

@contextmanager
def ocr_resource_manager():
    """Context manager for OCR resources cleanup"""
    temp_files = []
    resources = []
    try:
        yield temp_files, resources
    finally:
        # Clean up temporary files
        for temp_file in temp_files:
            try:
                if os.path.exists(temp_file):
                    os.unlink(temp_file)
            except Exception as e:
                print(f"Warning: Failed to delete temp file {temp_file}: {e}")
        
        # Clean up resources
        for resource in resources:
            try:
                if hasattr(resource, 'close'):
                    resource.close()
                del resource
            except Exception as e:
                print(f"Warning: Failed to cleanup resource: {e}")
        
        # Force garbage collection
        gc.collect()

def perform_ocr_on_page(page, page_num, dpi=150):
    """Perform OCR on a single page with proper resource management"""
    pix = None
    image = None
    text = ""
    
    try:
        # Create pixmap with specified DPI
        pix = page.get_pixmap(dpi=dpi)
        
        # Convert to PIL Image
        image_bytes = pix.tobytes("png")
        image = Image.open(io.BytesIO(image_bytes))
        
        # Perform OCR with optimized config
        ocr_config = '--oem 3 --psm 6'
        text = pytesseract.image_to_string(image, config=ocr_config)
        
    except Exception as ocr_error:
        print(f"OCR failed on page {page_num}: {ocr_error}")
        text = ""
    
    finally:
        # Explicit cleanup
        if image:
            image.close()
            del image
        if pix:
            del pix
        
        # Force garbage collection
        gc.collect()
    
    return text

def extract_text_pdf_with_ocr(file_path, word_limit=200):
    """Extract text from PDF with OCR fallback and proper memory management"""
    doc = None
    try:
        doc = fitz.open(file_path)
        word_list = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            
            if not text.strip():
                # OCR fallback with proper resource management
                text = perform_ocr_on_page(page, page_num, dpi=150)
                log_memory_usage(f"PDF OCR fallback page {page_num}")
            
            if text.strip():
                text = clean_text(text)
                words = text.split()
                word_list.extend(words)
                
                if len(word_list) >= word_limit:
                    break
            
            # Cleanup page reference
            page = None
            
            # Force garbage collection every 5 pages
            if page_num % 5 == 0:
                gc.collect()
        
        return ' '.join(word_list[:word_limit])
        
    except Exception as e:
        print(f"Error in extract_text_pdf_with_ocr: {e}")
        return ""
    
    finally:
        if doc:
            doc.close()
        gc.collect()

def extract_text_pdf_random(file_path, word_limit=1000):
    """Extract text from PDF with random page selection and proper memory management"""
    doc = None
    try:
        doc = fitz.open(file_path)
        words = []
        
        # First try to extract text normally
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            
            if not text.strip():
                break
                
            text = clean_text(text)
            words += text.split()
            
            if len(words) >= word_limit:
                break
            
            # Cleanup page reference
            page = None
            
            if page_num % 5 == 0:
                log_memory_usage(f"PDF random page {page_num}")
                gc.collect()
        
        # If no text found, try OCR on random pages
        if not words:
            total_pages = len(doc)
            
            # Select pages to OCR
            if total_pages >= 2:
                pages_to_ocr = [0, 1]
            else:
                pages_to_ocr = [0]
            
            # Add random pages (max 4 total)
            while len(set(pages_to_ocr)) < min(4, total_pages):
                pages_to_ocr.append(random.randint(0, total_pages - 1))
            
            # Remove duplicates and sort
            pages_to_ocr = sorted(list(set(pages_to_ocr)))
            
            # Process selected pages with OCR
            for page_num in pages_to_ocr:
                page = doc[page_num]
                text = perform_ocr_on_page(page, page_num, dpi=150)
                
                if text.strip():
                    text = clean_text(text)
                    words += text.split()
                    
                    if len(words) >= word_limit:
                        break
                
                # Cleanup page reference
                page = None
                
                log_memory_usage(f"PDF OCR fallback page {page_num}")
        
        return ' '.join(words[:word_limit]) if words else "No text found in the PDF."
        
    except Exception as e:
        print(f"Error in extract_text_pdf_random: {e}")
        return ""
    
    finally:
        if doc:
            doc.close()
        gc.collect()

def extract_text_docx(file_path, word_limit=200):
    """Extract text from DOCX with memory management"""
    doc = None
    try:
        doc = Document(file_path)
        words = []
        
        for para in doc.paragraphs:
            text = clean_text(para.text)
            words += text.split()
            if len(words) >= word_limit:
                break
        
        return ' '.join(words[:word_limit])
        
    except Exception as e:
        print(f"Error in extract_text_docx: {e}")
        return ""
    
    finally:
        if doc:
            del doc
        gc.collect()

def extract_text_pptx(file_path, word_limit=200):
    """Extract text from PPTX with memory management"""
    prs = None
    try:
        prs = Presentation(file_path)
        words = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = clean_text(shape.text)
                    words += text.split()
                    if len(words) >= word_limit:
                        break
            if len(words) >= word_limit:
                break
        
        return ' '.join(words[:word_limit])
        
    except Exception as e:
        print(f"Error in extract_text_pptx: {e}")
        return ""
    
    finally:
        if prs:
            del prs
        gc.collect()

def extract_text_txt(file_path, word_limit=200):
    """Extract text from TXT with memory management"""
    try:
        words = []
        with open(file_path, 'r', encoding='utf-8') as f:
            for line_num, line in enumerate(f):
                line = clean_text(line)
                words += line.split()
                if len(words) >= word_limit:
                    break
                
                # Log memory usage periodically
                if line_num % 100 == 0:
                    log_memory_usage(f"TXT line {line_num}")
        
        return ' '.join(words[:word_limit])
        
    except Exception as e:
        print(f"Error in extract_text_txt: {e}")
        return ""
    
    finally:
        gc.collect()

def extract_text_from_file(file_path, word_limit=1000):
    """Extract text from file with proper memory management"""
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        return ""
    
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        log_memory_usage(f"Before processing {ext} file")
        
        if ext == '.pdf':
            result = extract_text_pdf_random(file_path, word_limit)
        elif ext == '.docx':
            result = extract_text_docx(file_path, word_limit)
        elif ext == '.pptx':
            result = extract_text_pptx(file_path, word_limit)
        elif ext == '.txt':
            result = extract_text_txt(file_path, word_limit)
        else:
            raise ValueError("Unsupported file type: " + ext)
        
        log_memory_usage(f"After processing {ext} file")
        return result
        
    except Exception as e:
        print(f"Error in extract_text_from_file: {e}")
        return ""
    
    finally:
        # Force garbage collection after processing any file
        gc.collect()

def generate_description_from_text(text):
    if not text or len(text.split()) < 15:
        return "No sufficient content available to generate a description."

    try:
        client = genai.Client(api_key=os.getenv("GEMINI_API_KEY"))
        prompt = (
            "I will provide you with a small portion of text from some study material. "
            "Your task is to analyze the content and generate a short, clear description of what the full material is likely about. "
            "The description should be 2–3 sentences long, summarizing the overall topic and purpose, and suitable as a preview or caption for a notes-sharing platform. "
            "Use simple, professional language that helps students quickly understand what the content covers. "
            "Do not copy exact sentences—paraphrase instead. Avoid unnecessary details and stay focused on the main subject. "

            "If the input is empty or too short to understand the context, return this message: 'Not enough data to generate a description.' "

            "If the content seems to be from a question bank, return a description like: 'This is a set of question bank for [subject_name], containing important questions for practice and review.' "

            "If the input appears unrelated to academic content, or contains irrelevant or non-syllabus-based material, return the same message: 'Not enough data to generate a description or failed to extract text from the notes' "

            "Your only job is to generate a meaningful description or handle the input based on these instructions—do not do anything else."
        )

        response = client.models.generate_content(
            model="gemma-3-27b-it",
            contents=prompt + "\n\n" + text
        )
        return response.text
        
    except Exception as e:
        print(f"Error while generating summary: {e}")
        return ""
    
    finally:
        # Cleanup any potential resources
        gc.collect()

# Test function to verify Tesseract installation
def test_tesseract():
    """Test if Tesseract is properly installed and accessible"""
    img = None
    try:
        # Create a simple test image with text
        from PIL import Image, ImageDraw, ImageFont
        
        # Create a simple image with text
        img = Image.new('RGB', (200, 100), color='white')
        draw = ImageDraw.Draw(img)
        draw.text((10, 10), "Test OCR", fill='black')
        
        # Test OCR
        text = pytesseract.image_to_string(img)
        print(f"Tesseract test successful. Detected text: '{text.strip()}'")
        return True
        
    except Exception as e:
        print(f"Tesseract test failed: {e}")
        return False
    
    finally:
        if img:
            img.close()
        gc.collect()

def download_file_from_google_drive(file_id, file_name):
    """Downloads a file using its Google Drive file ID and saves it to temp/"""
    URL = "https://drive.google.com/uc?export=download"
    random_suffix = random.randint(100000, 999999)
    file_name = f"{random_suffix}_{file_name}"
    session = requests.Session()
    response = session.get(URL, params={'id': file_id}, stream=True)
    
    # Check for confirmation token for large files
    for key, value in response.cookies.items():
        if key.startswith('download_warning'):
            response = session.get(URL, params={'id': file_id, 'confirm': value}, stream=True)
            break

    temp_dir = TEMP_DIR
    os.makedirs(temp_dir, exist_ok=True)
    file_path = os.path.join(temp_dir, file_name)
    
    try:
        with open(file_path, "wb") as f:
            for chunk in response.iter_content(32768):
                if chunk:
                    f.write(chunk)
        print(f"[✓] File downloaded to {file_path}")
        return file_path
        
    except Exception as e:
        print("[✗] Error while saving the file to temp:", e)
        return None
    
    finally:
        if 'response' in locals():
            response.close()
        if 'session' in locals():
            session.close()
        gc.collect()
    
def start_worker_if_needed():
    global worker_active
    with worker_lock:
        if not worker_active:
            thread = Thread(target=generate_description_worker, args=(queue, semaphore), daemon=True)
            thread.start()
            worker_active = True

def process_description(note):
    """Process a single note with proper resource cleanup"""
    file_id = note.file_path
    file_name = note.filename
    temp_path = None

    try:
        log_memory_usage(f"Before processing {file_name}")
        
        temp_path = download_file_from_google_drive(file_id, file_name)
        if not temp_path:
            return

        text = extract_text_from_file(temp_path)
        description = generate_description_from_text(text)
        db.save_summary(drive_file_path=file_id, summary=description)

        print(f"[✓] Summary updated for file {file_name}")
        log_memory_usage(f"After processing {file_name}")
        
    except Exception as e:
        print(f"Error processing {file_name}: {e}")
    
    finally:
        # Always cleanup temp file
        if temp_path and os.path.exists(temp_path):
            try:
                os.remove(temp_path)
            except Exception as e:
                print(f"Failed to remove temp file {temp_path}: {e}")
        
        # Force garbage collection
        gc.collect()

def generate_description_worker(queue, semaphore):
    """Worker thread with enhanced memory management"""
    global worker_active
    processed_count = 0
    
    try:
        while not queue.empty():
            note = queue.get()
            with semaphore:
                try:
                    process_description(note)
                    processed_count += 1
                    
                    # Force garbage collection every 3 processed files
                    if processed_count % 3 == 0:
                        gc.collect()
                        log_memory_usage(f"After processing {processed_count} files")
                        
                except Exception as e:
                    print("Error in processing:", e)
            queue.task_done()
            
    except Exception as e:
        print(f"Worker thread error: {e}")
    
    finally:
        clear_temp_folder()
        gc.collect()
        with worker_lock:
            worker_active = False

def clear_temp_folder():
    """Clear temp folder with error handling"""
    if os.path.exists(TEMP_DIR):
        try:
            shutil.rmtree(TEMP_DIR)
            print(f"[✓] Cleared temp folder: {TEMP_DIR}")
        except Exception as e:
            print(f"[✗] Failed to clear temp folder: {e}")
    
    try:
        os.makedirs(TEMP_DIR, exist_ok=True)
    except Exception as e:
        print(f"[✗] Failed to create temp folder: {e}")

@app.route('/initialize_description_worker', methods=['POST', 'GET'])
def start_generating_description():
    try:
        null_notes = db.get_null_notes()
        if not null_notes:
            return jsonify({"message": "No notes found with empty descriptions."}), 404

        for note in null_notes:
            queue.put(note)

        start_worker_if_needed()

        return jsonify({"message": f"{len(null_notes)} jobs added to queue."}), 200
        
    except Exception as e:
        print(f"Error in start_generating_description: {e}")
        return jsonify({"message": "Failed to initialize description worker."}), 500

@app.route('/memory_status')
def memory_status():
    """Endpoint to check current memory usage"""
    process = psutil.Process(os.getpid())
    mem_info = process.memory_info()
    return jsonify({
        "rss_mb": round(mem_info.rss / (1024 * 1024), 2),
        "vms_mb": round(mem_info.vms / (1024 * 1024), 2),
        "queue_size": queue.qsize(),
        "worker_active": worker_active
    })

@app.route('/ping')
def ping():
    return "Summary service is alive", 200

if __name__ == "__main__":
    # Test Tesseract installation
    if not test_tesseract():
        print("Tesseract OCR is not properly installed or configured.")
    
    clear_temp_folder()
    
    # Start the Flask app
    app.run()